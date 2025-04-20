import os
import json
import re
import tempfile
import streamlit as st
from PyPDF2 import PdfReader
from groq import Groq
import speech_recognition as sr
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Constants
GROQ_API_KEY = ""
MODEL = "llama-3.3-70b-versatile"
USER_CREDENTIALS = {"teacher#7242": "log123", "student#1513": "log123"}

# ==============================
# INTENT CLASSIFICATION
# ==============================

def classify_intent(query, groq_client, model="llama-3.3-70b-versatile"):
    """
    Classifies if a query is related to PowerPoint generation or a normal query.
    
    Args:
        query (str): The user's input query
        groq_client: The initialized Groq client
        model (str): The model to use for classification
        
    Returns:
        str: Either "ppt_generation" or "normal_query"
    """
    # Simple rule-based checks for obvious PPT requests
    ppt_keywords = [
        "ppt", "powerpoint", "presentation", "slide", "slides", "deck", 
        "create a presentation", "make a presentation", "generate a presentation",
        "create a ppt", "make a ppt", "generate a ppt",
        "create slides", "make slides", "generate slides"
    ]
    
    # Check for obvious matches first to avoid unnecessary API calls
    query_lower = query.lower()
    for keyword in ppt_keywords:
        if keyword in query_lower:
            return "ppt_generation"
    
    # For more ambiguous queries, use the LLM for classification
    system_prompt = """
    You are an intent classifier. Determine if the user query is related to PowerPoint (PPT) 
    presentation generation or if it's a normal query.
    
    Return ONLY one of these exact strings:
    - "ppt_generation" - if the query is about creating, generating, or making a PowerPoint presentation
    - "normal_query" - for all other queries
    
    Be strict about classifying as "ppt_generation" - only classify it as such if the user is clearly 
    asking for a PowerPoint presentation to be created.
    """
    
    try:
        response = groq_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": query}
            ]
        )
        classification = response.choices[0].message.content.strip().lower()
        
        # Ensure we get one of our expected responses
        if "ppt_generation" in classification:
            return "ppt_generation"
        else:
            return "normal_query"
    except Exception:
        # Fall back to normal query in case of any errors
        return "normal_query"

# ==============================
# POWERPOINT GENERATION
# ==============================

def extract_json(text):
    """Extracts the first JSON object from a string."""
    match = re.search(r'\{.*\}', text, re.DOTALL)
    return match.group(0) if match else None

def generate_presentation_content(prompt, client, model="llama-3.3-70b-versatile"):
    """Generates structured presentation content using Groq LLM based on a prompt."""
    system_prompt = """
You are a presentation creation assistant. Create detailed, professional PowerPoint slides based on the user's prompt.
Return ONLY a JSON object with the following structure:
{
  "title": "<Presentation Title>",
  "slides": [
    {
      "title": "<Slide Title>",
      "type": "title_slide",
      "content": "<Main title content>",
      "subtitle": "<Optional subtitle>"
    },
    {
      "title": "<Slide Title>",
      "type": "bullet_points",
      "points": ["<Point 1>", "<Point 2>", "<Point 3>", "<Point 4>"]
    },
    {
      "title": "<Slide Title>",
      "type": "section_header",
      "content": "<Section title or transition slide content>"
    },
    {
      "title": "<Slide Title>",
      "type": "content_slide",
      "paragraphs": ["<Paragraph 1>", "<Paragraph 2>"]
    },
    {
      "title": "<Slide Title>",
      "type": "bullet_points",
      "points": ["<Point 1>", "<Point 2>", "<Point 3>", "<Point 4>"]
    },
    {
      "title": "<Slide Conclusion Title>",
      "type": "conclusion_slide",
      "content": "<Conclusion content>",
      "key_takeaway": "<Key takeaway message>"
    }
  ]
}
You MUST create EXACTLY 6 slides following this pattern:
1. Title slide with presentation title and subtitle
2. Introduction slide with 4 bullet points 
3. One key concept slide with a section header
4. Content slide with 2 paragraphs of information
5. Details slide with 4 key bullet points
6. Conclusion slide with summary and key takeaway

Ensure content is distributed evenly and symmetrically across all slides. Make the presentation visually balanced.
Do not include any text outside the JSON structure.
"""
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ]
        )
        output = response.choices[0].message.content.strip()
        json_str = extract_json(output)
        return json.loads(json_str) if json_str else {"error": "No valid JSON found."}
    except Exception as e:
        return {"error": str(e)}

def create_powerpoint(presentation_data, output_path):
    """Creates a PowerPoint presentation from the structured data with enhanced styling."""
    # Create a new presentation
    prs = Presentation()
    
    # Define some color schemes for a professional look
    color_schemes = [
        {"title": RGBColor(0, 112, 192), "text": RGBColor(0, 0, 0), "accent": RGBColor(79, 129, 189)},  # Blue theme
        {"title": RGBColor(192, 80, 77), "text": RGBColor(0, 0, 0), "accent": RGBColor(149, 55, 53)},   # Red theme
        {"title": RGBColor(0, 176, 80), "text": RGBColor(0, 0, 0), "accent": RGBColor(0, 112, 52)},     # Green theme
        {"title": RGBColor(112, 48, 160), "text": RGBColor(0, 0, 0), "accent": RGBColor(91, 15, 0)}     # Purple theme
    ]
    
    # Select a color scheme
    color_scheme = color_schemes[0]  # Default to blue theme
    
    # Ensure we have exactly 6 slides
    if len(presentation_data["slides"]) != 6:
        # Pad or truncate to exactly 6 slides
        if len(presentation_data["slides"]) < 6:
            # Add more slides if less than 6
            while len(presentation_data["slides"]) < 6:
                presentation_data["slides"].append({
                    "title": f"Additional Information",
                    "type": "bullet_points",
                    "points": ["Key point about the topic", "Supporting information", "Additional detail", "Final point"]
                })
        else:
            # Truncate to 6 slides if more
            presentation_data["slides"] = presentation_data["slides"][:6]
    
    # Create slides based on the presentation data
    for index, slide_data in enumerate(presentation_data["slides"]):
        # Select layout based on slide type
        if slide_data["type"] == "title_slide":
            layout = prs.slide_layouts[0]  # Title Slide layout
        elif slide_data["type"] == "section_header":
            layout = prs.slide_layouts[2]  # Section Header layout
        elif slide_data["type"] in ["bullet_points", "conclusion_slide"]:
            layout = prs.slide_layouts[1]  # Title and Content layout
        elif slide_data["type"] == "content_slide":
            layout = prs.slide_layouts[5]  # Blank layout
        else:
            layout = prs.slide_layouts[1]  # Default to Title and Content
        
        # Add slide
        slide = prs.slides.add_slide(layout)
        
        # Add title to all slides with styling
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]
            
            # Style the title
            title_frame = title_shape.text_frame
            title_frame.text = slide_data["title"]
            for paragraph in title_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(32)
                    run.font.color.rgb = color_scheme["title"]
        
        # Add content based on slide type with styling
        if slide_data["type"] == "title_slide":
            if len(slide.placeholders) > 1:  # If there's a subtitle placeholder
                subtitle = slide.placeholders[1]
                
                # Use subtitle field if available, otherwise use content
                subtitle_text = slide_data.get("subtitle", slide_data.get("content", ""))
                subtitle.text = subtitle_text
                
                # Style the subtitle
                for paragraph in subtitle.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.italic = True
        
        elif slide_data["type"] == "bullet_points":
            if len(slide.shapes.placeholders) > 1:  # Content placeholder
                content = slide.placeholders[1]
                tf = content.text_frame
                
                # Clear any existing paragraphs
                if tf.paragraphs:
                    tf.clear()
                
                # Add each bullet point with styling
                for point in slide_data["points"]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0  # Main bullet point
                    p.alignment = PP_ALIGN.LEFT
                    
                    # Style the bullet point text
                    for run in p.runs:
                        run.font.size = Pt(20)
                        run.font.color.rgb = color_scheme["text"]
        
        elif slide_data["type"] == "content_slide":
            # For content slides, add a text box with centered content
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            
            # Add paragraphs with styling
            for i, paragraph_text in enumerate(slide_data.get("paragraphs", [])):
                p = tf.add_paragraph()
                p.text = paragraph_text
                p.alignment = PP_ALIGN.LEFT
                
                # Add spacing between paragraphs
                if i > 0:
                    p.space_before = Pt(12)
                    
                # Style the paragraph text
                for run in p.runs:
                    run.font.size = Pt(20)
                    run.font.color.rgb = color_scheme["text"]
                
        elif slide_data["type"] == "section_header":
            # Section headers with styled content
            if "content" in slide_data and len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                content.text = slide_data["content"]
                
                # Style the section header content
                for paragraph in content.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(28)
                        run.font.bold = True
                        run.font.color.rgb = color_scheme["accent"]
        
        elif slide_data["type"] == "conclusion_slide":
            # Handle conclusion slides similarly to bullet points
            if len(slide.shapes.placeholders) > 1:
                content = slide.placeholders[1]
                tf = content.text_frame
                
                # Add main conclusion content
                p = tf.add_paragraph()
                p.text = slide_data.get("content", "")
                p.alignment = PP_ALIGN.LEFT
                
                # Style the content
                for run in p.runs:
                    run.font.size = Pt(22)
                    run.font.color.rgb = color_scheme["text"]
                
                # Add space before key takeaway
                if "key_takeaway" in slide_data:
                    p = tf.add_paragraph()
                    p.space_before = Pt(20)
                    p.text = ""  # Empty paragraph for spacing
                    
                    # Add key takeaway with special styling
                    p = tf.add_paragraph()
                    p.text = f"Key Takeaway: {slide_data['key_takeaway']}"
                    p.alignment = PP_ALIGN.CENTER
                    
                    # Style the key takeaway
                    for run in p.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
                        run.font.italic = True
                        run.font.color.rgb = color_scheme["accent"]
    
    # Save the presentation
    prs.save(output_path)
    return output_path

def process_ppt_request(user_input, groq_client):
    """Process a PPT generation request and return the presentation file."""
    # Extract topic from the query (improved pattern matching)
    topic_match = re.search(r'(?:generate|create|make|prepare)(?:\s+a)?\s+(?:ppt|powerpoint|presentation|slides)(?:\s+on|about|for|covering)?\s+(.*)', user_input.lower())
    topic = topic_match.group(1) if topic_match else user_input
    
    # Format the topic for better presentation title
    topic = topic.strip().capitalize()
    
    # Create a success message container
    success_container = st.empty()
    success_container.info("🤔 Analyzing your request and planning the presentation structure...")
    
    with st.spinner("🔄 Generating presentation content with exactly 6 slides..."):
        presentation_data = generate_presentation_content(topic, groq_client)
    
    if "error" in presentation_data:
        return f"❌ Error generating presentation: {presentation_data['error']}"
    
    # Update status
    success_container.info("📊 Designing a visually balanced PowerPoint presentation...")
    
    with st.spinner("✨ Creating a professionally styled PowerPoint presentation..."):
        output_file = f"{topic.replace(' ', '_')}_presentation.pptx"
        pptx_path = create_powerpoint(presentation_data, output_file)
        
        # Show presentation structure
        st.subheader("📑 Presentation Structure")
        st.write(f"**Title:** {presentation_data['title']}")
        
        # Display slide list
        for i, slide in enumerate(presentation_data["slides"]):
            st.write(f"**Slide {i+1}:** {slide['title']} ({slide['type']})")
        
        # Read the created file for download
        with open(pptx_path, "rb") as file:
            download_data = file.read()
            
        st.download_button(
            label="📄 Download Presentation",
            data=download_data,
            file_name=output_file,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"download_ppt_{topic}"
        )
        
        # Clear the info message
        success_container.empty()
        
        return f"✅ Your 6-slide presentation on '{topic}' is ready for download! Each slide has been professionally designed with symmetrical content distribution."

# ==============================
# CHAT FUNCTIONS
# ==============================

def process_chat_query(query, groq_client, model="llama-3.3-70b-versatile"):
    """Process a normal chat query and return the response."""
    try:
        messages = st.session_state.messages.copy()
        # Add user message if not already added
        if not messages or messages[-1]["role"] != "user" or messages[-1]["content"] != query:
            messages.append({"role": "user", "content": query})
            
        completion = groq_client.chat.completions.create(
            model=model,
            messages=messages
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"Error generating response: {str(e)}"

def generate_speech(text, lang="en"):
    """Generate speech from text and return the temporary file path."""
    try:
        tts = gTTS(text=text, lang=lang)
        temp_audio_file = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
        tts.save(temp_audio_file.name)
        return temp_audio_file.name
    except Exception as e:
        st.error(f"Audio error: {str(e)}")
        return None

# ==============================
# UI FUNCTIONS
# ==============================

def apply_styling():
    """Apply custom styling to the UI."""
    st.markdown("""
        <style>
        * { font-size: 16px !important; }
        .main-header { font-size: 24px !important; font-weight: bold; text-align: center; color: #4a4a4a; }
        .sub-header { font-size: 18px !important; text-align: center; color: #6e6e6e; margin: 10px 0; }
        .upload-section { background-color: #f8f9fa; padding: 20px; border-radius: 10px; margin-top: 15px; }
        .stButton button { width: 100%; font-size: 16px !important; }
        .login-container {
            max-width: 400px;
            margin: 0 auto;
            padding: 20px;
            border-radius: 10px;
            background-color: #f8f9fa;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        }
        .login-header {
            text-align: center;
            font-size: 24px !important;
            color: #4a4a4a;
            margin-bottom: 20px;
        }
        </style>
    """, unsafe_allow_html=True)

def login_page():
    """Display the login page and handle authentication."""
    # Create a centered container for the login form
    _, center_col, _ = st.columns([1, 2, 1])

    with center_col:
        st.markdown('<h2 class="login-header">🔐 Login</h2>', unsafe_allow_html=True)
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        if st.button("Login"):
            if username in USER_CREDENTIALS and USER_CREDENTIALS[username] == password:
                st.session_state.authenticated = True
                st.success("Login successful!")
            else:
                st.error("Invalid credentials")

def display_sidebar_tools():
    """Display and handle sidebar tools."""
    with st.sidebar:
        st.markdown('<p class="main-header">AI Assistant Tools</p>', unsafe_allow_html=True)

        task_options = [
            '📝 Auto-Grader',
            '🧾 Summarize Document',
            '🌐 Translate Content',
            '🎙️ Transcript Generation',
            '📊 Rubric Generation',
            '📊 PowerPoint Generation'
        ]
        selected_task = st.selectbox("Choose a task", task_options, label_visibility="collapsed")

        # File uploader
        pdf_file = st.file_uploader("Upload Document", type=["pdf"])

        # Help and information
        with st.expander("ℹ️ About This App", expanded=False):
            st.write("""
            This AI assistant helps with various educational tasks including 
            grading, document summarization, translation, transcript creation, 
            and PowerPoint generation.
            """)

def main():
    """Main application entry point."""
    # Page configuration
    st.set_page_config(page_title="AI Assistant", page_icon="🤖", layout="wide")
    
    # Apply custom styling
    apply_styling()
    
    # Initialize Groq client
    @st.cache_resource
    def get_groq_client():
        return Groq(api_key=GROQ_API_KEY)
    groq_client = get_groq_client()
    
    # Initialize session state
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "file_content" not in st.session_state:
        st.session_state.file_content = None
    
    # If not authenticated, show login page
    if not st.session_state.authenticated:
        # Add a header before the login form
        st.markdown('<h1 class="main-header">🤖 AI Assistant</h1>', unsafe_allow_html=True)
        st.markdown('<p class="sub-header">EDUCATIONAL TOOLS</p>', unsafe_allow_html=True)
        
        login_page()
        st.stop()
    
    # Display sidebar tools
    display_sidebar_tools()
    
    # Main area
    col1, col2, col3 = st.columns([1, 10, 1])
    with col2:
        st.markdown('<h1 class="main-header">🤖 AI Assistant</h1>', unsafe_allow_html=True)
        st.markdown('<p class="sub-header">EDUCATIONAL TOOLS</p>', unsafe_allow_html=True)
        
        # Chat history display
        chat_container = st.container()
        with chat_container:
            for msg in st.session_state.messages:
                with st.chat_message(msg["role"]):
                    st.markdown(msg["content"])
                    if msg["role"] == "assistant":
                        # Add audio for assistant messages
                        audio_path = generate_speech(msg["content"])
                        if audio_path:
                            st.audio(audio_path, format="audio/mp3")
        
        # Chat input
        chat_input_container = st.container()
        with chat_input_container:
            input_col, voice_col = st.columns([6, 1])
            with input_col:
                user_input = st.chat_input("Type your query here...")
            with voice_col:
                voice_button = st.button("🎙️", key="voice_input")
        
        # Voice input handling
        if voice_button:
            with st.spinner("Listening... Speak now"):
                try:
                    recognizer = sr.Recognizer()
                    with sr.Microphone() as source:
                        audio = recognizer.listen(source, timeout=5)
                        query = recognizer.recognize_google(audio)
                        
                        # Process the voice query
                        st.session_state.messages.append({"role": "user", "content": query})
                        with st.chat_message("user"):
                            st.markdown(query)
                        
                        # Classify the intent
                        intent = classify_intent(query, groq_client)
                        
                        with st.chat_message("assistant"):
                            if intent == "ppt_generation":
                                # Process as PPT generation request
                                reply = process_ppt_request(query, groq_client)
                                st.markdown(reply)
                            else:
                                # Process as normal query
                                with st.spinner("Generating response..."):
                                    reply = process_chat_query(query, groq_client)
                                    st.markdown(reply)
                            
                            # Generate audio response
                            audio_path = generate_speech(reply)
                            if audio_path:
                                st.audio(audio_path, format="audio/mp3")
                            
                            # Add assistant response to message history
                            st.session_state.messages.append({"role": "assistant", "content": reply})
                except Exception as e:
                    st.error(f"Speech recognition failed: {e}")
        
        # Text input handling
        if user_input:
            # Add user message to history
            st.session_state.messages.append({"role": "user", "content": user_input})
            with st.chat_message("user"):
                st.markdown(user_input)
            
            # Classify the intent using the intent classifier
            intent = classify_intent(user_input, groq_client)
            
            with st.chat_message("assistant"):
                if intent == "ppt_generation":
                    # Process as PPT generation request
                    reply = process_ppt_request(user_input, groq_client)
                    st.markdown(reply)
                    
                    # Generate audio response
                    audio_path = generate_speech(reply)
                    if audio_path:
                        st.audio(audio_path, format="audio/mp3")
                    
                    # Add assistant response to message history
                    st.session_state.messages.append({"role": "assistant", "content": reply})
                
                else:
                    # Process as normal query
                    with st.spinner("Thinking..."):
                        try:
                            reply = process_chat_query(user_input, groq_client)
                            st.markdown(reply)

                            # Generate audio response
                            audio_path = generate_speech(reply)
                            if audio_path:
                                st.audio(audio_path, format="audio/mp3")

                            # Add download button for response
                            st.download_button(
                                label="📥 Download Response", 
                                data=reply, 
                                file_name="ai_response.txt", 
                                mime="text/plain"
                            )
                            
                            # Add assistant response to message history
                            st.session_state.messages.append({"role": "assistant", "content": reply})
                        except Exception as e:
                            st.error(f"Response generation failed: {e}")

if __name__ == "__main__":
    main()
