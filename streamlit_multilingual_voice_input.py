import os
import json
import re
import tempfile
import streamlit as st
from PyPDF2 import PdfReader
from groq import Groq
import speech_recognition as sr
from gtts import gTTS
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Constants
GROQ_API_KEY = "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"
MODEL = "llama-3.3-70b-versatile"
USER_CREDENTIALS = {"teacher#7242": "log123", "student#1513": "log123"}

# Indian language support
INDIAN_LANGUAGES = {
    "English": {"code": "en-IN", "tts_code": "en", "display": "English"},
    "Hindi": {"code": "hi-IN", "tts_code": "hi", "display": "हिंदी"},
    "Bengali": {"code": "bn-IN", "tts_code": "bn", "display": "বাংলা"},
    "Tamil": {"code": "ta-IN", "tts_code": "ta", "display": "தமிழ்"},
    "Telugu": {"code": "te-IN", "tts_code": "te", "display": "తెలుగు"},
    "Marathi": {"code": "mr-IN", "tts_code": "mr", "display": "मराठी"},
    "Gujarati": {"code": "gu-IN", "tts_code": "gu", "display": "ગુજરાતી"},
    "Kannada": {"code": "kn-IN", "tts_code": "kn", "display": "ಕನ್ನಡ"},
    "Malayalam": {"code": "ml-IN", "tts_code": "ml", "display": "മലയാളം"},
    "Punjabi": {"code": "pa-IN", "tts_code": "pa", "display": "ਪੰਜਾਬੀ"},
    "Urdu": {"code": "ur-IN", "tts_code": "ur", "display": "اردو"}
}

# ==============================
# INTENT CLASSIFICATION
# ==============================

def classify_intent(query, groq_client, model="llama-3.3-70b-versatile"):
    """
    Classifies if a query is related to PowerPoint generation or a normal query.
    
    Args:
        query (str): The user's input query
        groq_client: The initialized Groq client
        model (str): The model to use for classification
        
    Returns:
        str: Either "ppt_generation" or "normal_query"
    """
    # Simple rule-based checks for obvious PPT requests
    ppt_keywords = [
        "ppt", "powerpoint", "presentation", "slide", "slides", "deck", 
        "create a presentation", "make a presentation", "generate a presentation",
        "create a ppt", "make a ppt", "generate a ppt",
        "create slides", "make slides", "generate slides",
        # Hindi keywords
        "प्रेजेंटेशन", "पावरपॉइंट", "पीपीटी", "स्लाइड", 
        # Bengali keywords
        "পাওয়ার পয়েন্ট", "প্রেজেন্টেশন", "স্লাইড",
        # Tamil keywords
        "விளக்கக்காட்சி", "பவர்பாயிண்ட்", "ஸ்லைடு",
        # Telugu keywords
        "ప్రజెంటేషన్", "పవర్ పాయింట్", "స్లైడ్"
    ]
    
    # Check for obvious matches first to avoid unnecessary API calls
    query_lower = query.lower()
    for keyword in ppt_keywords:
        if keyword.lower() in query_lower:
            return "ppt_generation"
    
    # For more ambiguous queries, use the LLM for classification
    system_prompt = """
    You are an intent classifier. Determine if the user query is related to PowerPoint (PPT) 
    presentation generation or if it's a normal query. The query might be in any Indian language.
    
    Return ONLY one of these exact strings:
    - "ppt_generation" - if the query is about creating, generating, or making a PowerPoint presentation
    - "normal_query" - for all other queries
    
    Be strict about classifying as "ppt_generation" - only classify it as such if the user is clearly 
    asking for a PowerPoint presentation to be created.
    """
    
    try:
        response = groq_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": query}
            ]
        )
        classification = response.choices[0].message.content.strip().lower()
        
        # Ensure we get one of our expected responses
        if "ppt_generation" in classification:
            return "ppt_generation"
        else:
            return "normal_query"
    except Exception:
        # Fall back to normal query in case of any errors
        return "normal_query"

# ==============================
# POWERPOINT GENERATION
# ==============================

def extract_json(text):
    """Extracts the first JSON object from a string."""
    match = re.search(r'\{.*\}', text, re.DOTALL)
    return match.group(0) if match else None

def generate_presentation_content(prompt, client, model="llama-3.3-70b-versatile"):
    """Generates structured presentation content using Groq LLM based on a prompt."""
    system_prompt = """
You are a presentation creation assistant. Create detailed, professional PowerPoint slides based on the user's prompt.
The prompt may be in any language, but your response should be in the same language as the prompt.

Return ONLY a JSON object with the following structure:
{
  "title": "<Presentation Title>",
  "slides": [
    {
      "title": "<Slide Title>",
      "type": "title_slide",
      "content": "<Main title content>",
      "subtitle": "<Optional subtitle>"
    },
    {
      "title": "<Slide Title>",
      "type": "bullet_points",
      "points": ["<Point 1>", "<Point 2>", "<Point 3>", "<Point 4>"]
    },
    {
      "title": "<Slide Title>",
      "type": "section_header",
      "content": "<Section title or transition slide content>"
    },
    {
      "title": "<Slide Title>",
      "type": "content_slide",
      "paragraphs": ["<Paragraph 1>", "<Paragraph 2>"]
    },
    {
      "title": "<Slide Title>",
      "type": "bullet_points",
      "points": ["<Point 1>", "<Point 2>", "<Point 3>", "<Point 4>"]
    },
    {
      "title": "<Slide Conclusion Title>",
      "type": "conclusion_slide",
      "content": "<Conclusion content>",
      "key_takeaway": "<Key takeaway message>"
    }
  ]
}
You MUST create EXACTLY 6 slides following this pattern:
1. Title slide with presentation title and subtitle
2. Introduction slide with 4 bullet points 
3. One key concept slide with a section header
4. Content slide with 2 paragraphs of information
5. Details slide with 4 key bullet points
6. Conclusion slide with summary and key takeaway

Ensure content is distributed evenly and symmetrically across all slides. Make the presentation visually balanced.
Do not include any text outside the JSON structure.
"""
    try:
        response = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": prompt}
            ]
        )
        output = response.choices[0].message.content.strip()
        json_str = extract_json(output)
        return json.loads(json_str) if json_str else {"error": "No valid JSON found."}
    except Exception as e:
        return {"error": str(e)}

def create_powerpoint(presentation_data, output_path):
    """Creates a PowerPoint presentation from the structured data with enhanced styling."""
    # Create a new presentation
    prs = Presentation()
    
    # Define some color schemes for a professional look
    color_schemes = [
        {"title": RGBColor(0, 112, 192), "text": RGBColor(0, 0, 0), "accent": RGBColor(79, 129, 189)},  # Blue theme
        {"title": RGBColor(192, 80, 77), "text": RGBColor(0, 0, 0), "accent": RGBColor(149, 55, 53)},   # Red theme
        {"title": RGBColor(0, 176, 80), "text": RGBColor(0, 0, 0), "accent": RGBColor(0, 112, 52)},     # Green theme
        {"title": RGBColor(112, 48, 160), "text": RGBColor(0, 0, 0), "accent": RGBColor(91, 15, 0)}     # Purple theme
    ]
    
    # Select a color scheme
    color_scheme = color_schemes[0]  # Default to blue theme
    
    # Ensure we have exactly 6 slides
    if len(presentation_data["slides"]) != 6:
        # Pad or truncate to exactly 6 slides
        if len(presentation_data["slides"]) < 6:
            # Add more slides if less than 6
            while len(presentation_data["slides"]) < 6:
                presentation_data["slides"].append({
                    "title": f"Additional Information",
                    "type": "bullet_points",
                    "points": ["Key point about the topic", "Supporting information", "Additional detail", "Final point"]
                })
        else:
            # Truncate to 6 slides if more
            presentation_data["slides"] = presentation_data["slides"][:6]
    
    # Create slides based on the presentation data
    for index, slide_data in enumerate(presentation_data["slides"]):
        # Select layout based on slide type
        if slide_data["type"] == "title_slide":
            layout = prs.slide_layouts[0]  # Title Slide layout
        elif slide_data["type"] == "section_header":
            layout = prs.slide_layouts[2]  # Section Header layout
        elif slide_data["type"] in ["bullet_points", "conclusion_slide"]:
            layout = prs.slide_layouts[1]  # Title and Content layout
        elif slide_data["type"] == "content_slide":
            layout = prs.slide_layouts[5]  # Blank layout
        else:
            layout = prs.slide_layouts[1]  # Default to Title and Content
        
        # Add slide
        slide = prs.slides.add_slide(layout)
        
        # Add title to all slides with styling
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            title_shape = slide.shapes.title
            title_shape.text = slide_data["title"]
            
            # Style the title
            title_frame = title_shape.text_frame
            title_frame.text = slide_data["title"]
            for paragraph in title_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.bold = True
                    run.font.size = Pt(32)
                    run.font.color.rgb = color_scheme["title"]
        
        # Add content based on slide type with styling
        if slide_data["type"] == "title_slide":
            if len(slide.placeholders) > 1:  # If there's a subtitle placeholder
                subtitle = slide.placeholders[1]
                
                # Use subtitle field if available, otherwise use content
                subtitle_text = slide_data.get("subtitle", slide_data.get("content", ""))
                subtitle.text = subtitle_text
                
                # Style the subtitle
                for paragraph in subtitle.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(24)
                        run.font.italic = True
        
        elif slide_data["type"] == "bullet_points":
            if len(slide.shapes.placeholders) > 1:  # Content placeholder
                content = slide.placeholders[1]
                tf = content.text_frame
                
                # Clear any existing paragraphs
                if tf.paragraphs:
                    tf.clear()
                
                # Add each bullet point with styling
                for point in slide_data["points"]:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0  # Main bullet point
                    p.alignment = PP_ALIGN.LEFT
                    
                    # Style the bullet point text
                    for run in p.runs:
                        run.font.size = Pt(20)
                        run.font.color.rgb = color_scheme["text"]
        
        elif slide_data["type"] == "content_slide":
            # For content slides, add a text box with centered content
            left = Inches(1)
            top = Inches(2)
            width = Inches(8)
            height = Inches(5)
            
            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.word_wrap = True
            
            # Add paragraphs with styling
            for i, paragraph_text in enumerate(slide_data.get("paragraphs", [])):
                p = tf.add_paragraph()
                p.text = paragraph_text
                p.alignment = PP_ALIGN.LEFT
                
                # Add spacing between paragraphs
                if i > 0:
                    p.space_before = Pt(12)
                    
                # Style the paragraph text
                for run in p.runs:
                    run.font.size = Pt(20)
                    run.font.color.rgb = color_scheme["text"]
                
        elif slide_data["type"] == "section_header":
            # Section headers with styled content
            if "content" in slide_data and len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                content.text = slide_data["content"]
                
                # Style the section header content
                for paragraph in content.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    for run in paragraph.runs:
                        run.font.size = Pt(28)
                        run.font.bold = True
                        run.font.color.rgb = color_scheme["accent"]
        
        elif slide_data["type"] == "conclusion_slide":
            # Handle conclusion slides similarly to bullet points
            if len(slide.shapes.placeholders) > 1:
                content = slide.placeholders[1]
                tf = content.text_frame
                
                # Add main conclusion content
                p = tf.add_paragraph()
                p.text = slide_data.get("content", "")
                p.alignment = PP_ALIGN.LEFT
                
                # Style the content
                for run in p.runs:
                    run.font.size = Pt(22)
                    run.font.color.rgb = color_scheme["text"]
                
                # Add space before key takeaway
                if "key_takeaway" in slide_data:
                    p = tf.add_paragraph()
                    p.space_before = Pt(20)
                    p.text = ""  # Empty paragraph for spacing
                    
                    # Add key takeaway with special styling
                    p = tf.add_paragraph()
                    p.text = f"Key Takeaway: {slide_data['key_takeaway']}"
                    p.alignment = PP_ALIGN.CENTER
                    
                    # Style the key takeaway
                    for run in p.runs:
                        run.font.size = Pt(24)
                        run.font.bold = True
                        run.font.italic = True
                        run.font.color.rgb = color_scheme["accent"]
    
    # Save the presentation
    prs.save(output_path)
    return output_path

def process_ppt_request(user_input, groq_client):
    """Process a PPT generation request and return the presentation file."""
    # Extract topic from the query (improved pattern matching)
    topic_match = re.search(r'(?:generate|create|make|prepare)(?:\s+a)?\s+(?:ppt|powerpoint|presentation|slides)(?:\s+on|about|for|covering)?\s+(.*)', user_input.lower())
    topic = topic_match.group(1) if topic_match else user_input
    
    # Format the topic for better presentation title
    topic = topic.strip()
    if topic[0].isalpha():
        topic = topic[0].upper() + topic[1:]
    
    # Create a success message container
    success_container = st.empty()
    success_container.info("🤔 Analyzing your request and planning the presentation structure...")
    
    with st.spinner("🔄 Generating presentation content with exactly 6 slides..."):
        presentation_data = generate_presentation_content(topic, groq_client)
    
    if "error" in presentation_data:
        return f"❌ Error generating presentation: {presentation_data['error']}"
    
    # Update status
    success_container.info("📊 Designing a visually balanced PowerPoint presentation...")
    
    with st.spinner("✨ Creating a professionally styled PowerPoint presentation..."):
        output_file = f"{topic.replace(' ', '_')}_presentation.pptx"
        pptx_path = create_powerpoint(presentation_data, output_file)
        
        # Show presentation structure
        st.subheader("📑 Presentation Structure")
        st.write(f"**Title:** {presentation_data['title']}")
        
        # Display slide list
        for i, slide in enumerate(presentation_data["slides"]):
            st.write(f"**Slide {i+1}:** {slide['title']} ({slide['type']})")
        
        # Read the created file for download
        with open(pptx_path, "rb") as file:
            download_data = file.read()
            
        st.download_button(
            label="📄 Download Presentation",
            data=download_data,
            file_name=output_file,
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            key=f"download_ppt_{topic}"
        )
        
        # Clear the info message
        success_container.empty()
        
        return f"✅ Your 6-slide presentation on '{topic}' is ready for download! Each slide has been professionally designed with symmetrical content distribution."

# ==============================
# LANGUAGE DETECTION
# ==============================

def detect_language(text, groq_client, model="llama-3.3-70b-versatile"):
    """
    Detect the language of the input text using the LLM.
    
    Args:
        text (str): The input text to analyze
        groq_client: The initialized Groq client
        model (str): The model to use
        
    Returns:
        str: The detected language code or "en" if detection fails
    """
    system_prompt = """
    You are a language detection expert. Given an input text, determine which language it is written in.
    Focus specifically on Indian languages including Hindi, Bengali, Tamil, Telugu, Marathi, Gujarati, 
    Kannada, Malayalam, Punjabi, and Urdu, as well as English.
    
    Return ONLY the language name as a single word, for example:
    - "English"
    - "Hindi"
    - "Bengali"
    - "Tamil"
    - "Telugu"
    - "Marathi"
    - "Gujarati"
    - "Kannada"
    - "Malayalam"
    - "Punjabi"
    - "Urdu"
    
    If you're not sure, default to "English".
    """
    
    try:
        response = groq_client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": text}
            ]
        )
        detected_language = response.choices[0].message.content.strip()
        
        # Clean up the response to get just the language name
        detected_language = detected_language.split()[0].strip('.,:"\'')
        
        # Check if the detected language is in our supported languages
        if detected_language in INDIAN_LANGUAGES:
            return detected_language
        # Try case-insensitive match
        for lang in INDIAN_LANGUAGES:
            if detected_language.lower() == lang.lower():
                return lang
        
        # Default to English if no match found
        return "English"
    except Exception as e:
        st.error(f"Language detection error: {str(e)}")
        return "English"  # Default to English if detection fails

# ==============================
# CHAT FUNCTIONS
# ==============================

def process_chat_query(query, groq_client, model="llama-3.3-70b-versatile"):
    """Process a normal chat query and return the response."""
    try:
        # Create a copy of the messages but remove the 'language' field from each message
        # as the Groq API doesn't support this field
        api_messages = []
        for msg in st.session_state.messages:
            # Create a clean copy with only the required fields for the API
            api_msg = {
                "role": msg["role"],
                "content": msg["content"]
            }
            api_messages.append(api_msg)
            
        # Add user message if not already added
        if not api_messages or api_messages[-1]["role"] != "user" or api_messages[-1]["content"] != query:
            api_messages.append({"role": "user", "content": query})
            
        completion = groq_client.chat.completions.create(
            model=model,
            messages=api_messages
        )
        return completion.choices[0].message.content
    except Exception as e:
        return f"Error generating response: {str(e)}"

def generate_speech(text, lang_code="en"):
    """Generate speech from text and return the temporary file path."""
    try:
        tts = gTTS(text=text, lang=lang_code)
        temp_audio_file = tempfile.NamedTemporaryFile(delete=False, suffix=".mp3")
        tts.save(temp_audio_file.name)
        return temp_audio_file.name
    except Exception as e:
        st.error(f"Audio error: {str(e)}")
        return None

# ==============================
# UI FUNCTIONS
# ==============================

def apply_styling():
    """Apply custom styling to the UI."""
    st.markdown("""
        <style>
        * { font-size: 16px !important; }
        .main-header { font-size: 24px !important; font-weight: bold; text-align: center; color: #4a4a4a; }
        .sub-header { font-size: 18px !important; text-align: center; color: #6e6e6e; margin: 10px 0; }
        .upload-section { background-color: #f8f9fa; padding: 20px; border-radius: 10px; margin-top: 15px; }
        .stButton button { width: 100%; font-size: 16px !important; }
        .login-container {
            max-width: 400px;
            margin: 0 auto;
            padding: 20px;
            border-radius: 10px;
            background-color: #f8f9fa;
            box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        }
        .login-header {
            text-align: center;
            font-size: 24px !important;
            color: #4a4a4a;
            margin-bottom: 20px;
        }
        .language-selector {
            padding: 10px;
            background-color: #f0f2f6;
            border-radius: 5px;
            margin-bottom: 15px;
        }
        </style>
    """, unsafe_allow_html=True)

def login_page():
    """Display the login page and handle authentication."""
    # Create a centered container for the login form
    _, center_col, _ = st.columns([1, 2, 1])

    with center_col:
        st.markdown('<h2 class="login-header">🔐 Login</h2>', unsafe_allow_html=True)
        username = st.text_input("Username")
        password = st.text_input("Password", type="password")
        if st.button("Login"):
            if username in USER_CREDENTIALS and USER_CREDENTIALS[username] == password:
                st.session_state.authenticated = True
                # Initialize language preference
                st.session_state.language = "English"
                st.success("Login successful!")
            else:
                st.error("Invalid credentials")

def display_sidebar_tools():
    """Display and handle sidebar tools."""
    with st.sidebar:
        st.markdown('<p class="main-header">AI Assistant Tools</p>', unsafe_allow_html=True)
        
        # Language selector
        st.markdown('<div class="language-selector">', unsafe_allow_html=True)
        st.subheader("🌐 Select Language")
        
        # Display language options with their native names
        language_options = [f"{lang_info['display']} ({lang})" for lang, lang_info in INDIAN_LANGUAGES.items()]
        selected_language_display = st.selectbox(
            "Choose your preferred language",
            language_options,
            index=language_options.index(f"{INDIAN_LANGUAGES[st.session_state.language]['display']} ({st.session_state.language})"),
            label_visibility="collapsed"
        )
        
        # Extract language key from the display format
        selected_language = selected_language_display.split("(")[1].split(")")[0]
        
        # Update session state if language changed
        if selected_language != st.session_state.language:
            st.session_state.language = selected_language
            st.rerun()
            
        st.markdown('</div>', unsafe_allow_html=True)

        # Tools
        st.subheader("🛠️ Tools")
        task_options = [
            '📝 Auto-Grader',
            '🧾 Summarize Document',
            '🌐 Translate Content',
            '🎙️ Transcript Generation',
            '📊 Rubric Generation',
            '📊 PowerPoint Generation'
        ]
        selected_task = st.selectbox("Choose a task", task_options, label_visibility="collapsed")

        # File uploader
        st.subheader("📄 Document Upload")
        pdf_file = st.file_uploader("Upload Document", type=["pdf"])

        # Help and information
        with st.expander("ℹ️ About This App", expanded=False):
            st.write("""
            This multilingual AI assistant helps with various educational tasks including 
            grading, document summarization, translation, transcript creation, 
            and PowerPoint generation. It supports multiple Indian languages.
            """)

def main():
    """Main application entry point."""
    # Page configuration
    st.set_page_config(page_title="Multilingual AI Assistant", page_icon="🌐", layout="wide")
    
    # Apply custom styling
    apply_styling()
    
    # Initialize Groq client
    @st.cache_resource
    def get_groq_client():
        return Groq(api_key=GROQ_API_KEY)
    groq_client = get_groq_client()
    
    # Initialize session state
    if "authenticated" not in st.session_state:
        st.session_state.authenticated = False
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "file_content" not in st.session_state:
        st.session_state.file_content = None
    if "language" not in st.session_state:
        st.session_state.language = "English"
    
    # If not authenticated, show login page
    if not st.session_state.authenticated:
        # Add a header before the login form
        st.markdown('<h1 class="main-header">🌐 Multilingual AI Assistant</h1>', unsafe_allow_html=True)
        st.markdown('<p class="sub-header">SUPPORTING INDIAN LANGUAGES</p>', unsafe_allow_html=True)
        
        login_page()
        st.stop()
    
    # Display sidebar tools
    display_sidebar_tools()
    
    # Get current language codes for speech recognition and TTS
    current_lang = st.session_state.language
    speech_lang_code = INDIAN_LANGUAGES[current_lang]["code"]
    tts_lang_code = INDIAN_LANGUAGES[current_lang]["tts_code"]
    
    # Main area
    col1, col2, col3 = st.columns([1, 10, 1])
    with col2:
        st.markdown(f'<h1 class="main-header">🌐 Multilingual AI Assistant</h1>', unsafe_allow_html=True)
        st.markdown(f'<p class="sub-header">Currently using: {INDIAN_LANGUAGES[current_lang]["display"]}</p>', unsafe_allow_html=True)
        
        # Chat history display
        chat_container = st.container()
        with chat_container:
            for msg in st.session_state.messages:
                with st.chat_message(msg["role"]):
                    st.markdown(msg["content"])
                    if msg["role"] == "assistant":
                        # Add audio for assistant messages
                        # Detect the language if not specified
                        msg_lang = msg.get("language", current_lang)
                        audio_path = generate_speech(
                            msg["content"],
                            INDIAN_LANGUAGES[msg_lang]["tts_code"]
                        )
                        if audio_path:
                            st.audio(audio_path, format="audio/mp3")
        
        # Chat input
        chat_input_container = st.container()
        with chat_input_container:
            input_col, voice_col = st.columns([6, 1])
            with input_col:
                placeholder_text = {
                    "English": "Type your query here...",
                    "Hindi": "अपना प्रश्न यहां लिखें...",
                    "Bengali": "আপনার প্রশ্ন এখানে টাইপ করুন...",
                    "Tamil": "உங்கள் கேள்வியை இங்கே தட்டச்சு செய்யவும்...",
                    "Telugu": "మీ ప్రశ్నను ఇక్కడ టైప్ చేయండి...",
                    "Marathi": "आपला प्रश्न येथे टाइप करा...",
                    "Gujarati": "તમારો પ્રશ્ન અહીં ટાઇપ કરો...",
                    "Kannada": "ನಿಮ್ಮ ಪ್ರಶ್ನೆಯನ್ನು ಇಲ್ಲಿ ಟೈಪ್ ಮಾಡಿ...",
                    "Malayalam": "നിങ്ങളുടെ ചോദ്യം ഇവിടെ ടൈപ്പ് ചെയ്യുക...",
                    "Punjabi": "ਆਪਣਾ ਸਵਾਲ ਇੱਥੇ ਟਾਈਪ ਕਰੋ...",
                    "Urdu": "اپنا سوال یہاں ٹائپ کریں..."
                }
                
                user_input = st.chat_input(
                    placeholder_text.get(current_lang, "Type your query here...")
                )
            with voice_col:
                voice_button_text = "🎙️" if current_lang == "English" else "🎙️"
                voice_button = st.button(voice_button_text, key="voice_input")
        
        # Voice input handling
        if voice_button:
            listening_text = {
                "English": "Listening... Speak now",
                "Hindi": "सुन रहा हूँ... अब बोलें",
                "Bengali": "শুনছি... এখন বলুন",
                "Tamil": "கேட்கிறேன்... இப்போது பேசுங்கள்",
                "Telugu": "వింటున్నాను... ఇప్పుడు మాట్లాడండి",
                "Marathi": "ऐकत आहे... आता बोला",
                "Gujarati": "સાંભળી રહ્યો છું... હવે બોલો",
                "Kannada": "ಕೇಳುತ್ತಿದ್ದೇನೆ... ಈಗ ಮಾತನಾಡಿ",
                "Malayalam": "കേൾക്കുന്നു... ഇപ്പോൾ സംസാരിക്കൂ",
                "Punjabi": "ਸੁਣ ਰਿਹਾ ਹਾਂ... ਹੁਣ ਬੋਲੋ",
                "Urdu": "سن رہا ہوں... اب بولیں"
            }
            
            with st.spinner(listening_text.get(current_lang, "Listening... Speak now")):
                try:
                    recognizer = sr.Recognizer()
                    with sr.Microphone() as source:
                        audio = recognizer.listen(source, timeout=5)
                        query = recognizer.recognize_google(audio, language=speech_lang_code)
                        
                        # Process the voice query
                        st.session_state.messages.append({"role": "user", "content": query, "language": current_lang})
                        with st.chat_message("user"):
                            st.markdown(query)
                        
                        # Classify the intent
                        intent = classify_intent(query, groq_client)
                        
                        with st.chat_message("assistant"):
                            if intent == "ppt_generation":
                                # Process as PPT generation request
                                reply = process_ppt_request(query, groq_client)
                                st.markdown(reply)
                            else:
                                # Process as normal query
                                with st.spinner("Generating response..."):
                                    reply = process_chat_query(query, groq_client)
                                    st.markdown(reply)
                            
                            # Generate audio response
                            audio_path = generate_speech(reply, tts_lang_code)
                            if audio_path:
                                st.audio(audio_path, format="audio/mp3")
                            
                            # Add assistant response to message history
                            st.session_state.messages.append({"role": "assistant", "content": reply, "language": current_lang})
                except Exception as e:
                    error_msg = {
                        "English": f"Speech recognition failed: {e}",
                        "Hindi": f"आवाज़ पहचानने में विफल: {e}",
                        "Bengali": f"কণ্ঠস্বর স্বীকৃতি ব্যর্থ হয়েছে: {e}",
                        "Tamil": f"பேச்சு அங்கீகாரம் தோல்வியடைந்தது: {e}",
                        "Telugu": f"స్పీచ్ రికగ్నిషన్ విఫలమైంది: {e}",
                        "Marathi": f"स्पीच रेकग्निशन अयशस्वी: {e}",
                        "Gujarati": f"સ્પીચ રેકગ્નિશન નિષ્ફળ: {e}",
                        "Kannada": f"ಧ್ವನಿ ಗುರುತಿಸುವಿಕೆ ವಿಫಲವಾಗಿದೆ: {e}",
                        "Malayalam": f"സ്പീച്ച് റെക്കഗ്നിഷൻ പരാജയപ്പെട്ടു: {e}",
                        "Punjabi": f"ਆਵਾਜ਼ ਪਛਾਣ ਅਸਫਲ: {e}",
                        "Urdu": f"تقریر کی شناخت ناکام ہو گئی: {e}"
                    }
                    st.error(error_msg.get(current_lang, f"Speech recognition failed: {e}"))
        
        # Text input handling
        if user_input:
            # Detect language if auto-detection is enabled
            detected_lang = detect_language(user_input, groq_client)
            
            # Add user message to history
            st.session_state.messages.append({"role": "user", "content": user_input, "language": detected_lang})
            with st.chat_message("user"):
                st.markdown(user_input)
            
            # Classify the intent using the intent classifier
            intent = classify_intent(user_input, groq_client)
            
            with st.chat_message("assistant"):
                if intent == "ppt_generation":
                    # Process as PPT generation request
                    reply = process_ppt_request(user_input, groq_client)
                    st.markdown(reply)
                    
                    # Generate audio response in detected language
                    audio_path = generate_speech(reply, INDIAN_LANGUAGES[detected_lang]["tts_code"])
                    if audio_path:
                        st.audio(audio_path, format="audio/mp3")
                    
                    # Add assistant response to message history
                    st.session_state.messages.append({"role": "assistant", "content": reply, "language": detected_lang})
                
                else:
                    # Process as normal query
                    processing_text = {
                        "English": "Thinking...",
                        "Hindi": "सोच रहा हूँ...",
                        "Bengali": "চিন্তা করছি...",
                        "Tamil": "சிந்திக்கிறேன்...",
                        "Telugu": "ఆలోచిస్తున్నాను...",
                        "Marathi": "विचार करत आहे...",
                        "Gujarati": "વિચારી રહ્યો છું...",
                        "Kannada": "ಯೋಚಿಸುತ್ತಿದ್ದೇನೆ...",
                        "Malayalam": "ചിന്തിക്കുന്നു...",
                        "Punjabi": "ਸੋਚ ਰਿਹਾ ਹਾਂ...",
                        "Urdu": "سوچ رہا ہوں..."
                    }
                    
                    with st.spinner(processing_text.get(detected_lang, "Thinking...")):
                        try:
                            reply = process_chat_query(user_input, groq_client)
                            st.markdown(reply)

                            # Generate audio response in detected language
                            audio_path = generate_speech(reply, INDIAN_LANGUAGES[detected_lang]["tts_code"])
                            if audio_path:
                                st.audio(audio_path, format="audio/mp3")

                            # Add download button for response
                            download_text = {
                                "English": "📥 Download Response",
                                "Hindi": "📥 जवाब डाउनलोड करें",
                                "Bengali": "📥 প্রতিক্রিয়া ডাউনলোড করুন",
                                "Tamil": "📥 பதிலைப் பதிவிறக்கவும்",
                                "Telugu": "📥 ప్రతిస్పందనను డౌన్‌లోడ్ చేయండి",
                                "Marathi": "📥 प्रतिसाद डाउनलोड करा",
                                "Gujarati": "📥 જવાબ ડાઉનલોડ કરો",
                                "Kannada": "📥 ಪ್ರತಿಕ್ರಿಯೆಯನ್ನು ಡೌನ್‌ಲೋಡ್ ಮಾಡಿ",
                                "Malayalam": "📥 മറുപടി ഡൗൺലോഡ് ചെയ്യുക",
                                "Punjabi": "📥 ਜਵਾਬ ਡਾਊਨਲੋਡ ਕਰੋ",
                                "Urdu": "📥 جواب ڈاؤن لوڈ کریں"
                            }
                            
                            st.download_button(
                                label=download_text.get(detected_lang, "📥 Download Response"), 
                                data=reply, 
                                file_name=f"ai_response_{detected_lang.lower()}.txt", 
                                mime="text/plain"
                            )
                            
                            # Add assistant response to message history
                            st.session_state.messages.append({"role": "assistant", "content": reply, "language": detected_lang})
                        except Exception as e:
                            error_msg = {
                                "English": f"Response generation failed: {e}",
                                "Hindi": f"जवाब उत्पन्न करने में विफल: {e}",
                                "Bengali": f"প্রতিক্রিয়া জেনারেশন ব্যর্থ হয়েছে: {e}",
                                "Tamil": f"பதில் உருவாக்கம் தோல்வியடைந்தது: {e}",
                                "Telugu": f"ప్రతిస్పందన జనరేషన్ విఫలమైంది: {e}",
                                "Marathi": f"प्रतिसाद जनरेशन अयशस्वी: {e}",
                                "Gujarati": f"જવાબ જનરેશન નિષ્ફળ: {e}",
                                "Kannada": f"ಪ್ರತಿಕ್ರಿಯೆ ರಚನೆ ವಿಫಲವಾಗಿದೆ: {e}",
                                "Malayalam": f"പ്രതികരണം ജനറേഷൻ പരാജയപ്പെട്ടു: {e}",
                                "Punjabi": f"ਜਵਾਬ ਜੇਨਰੇਸ਼ਨ ਅਸਫਲ: {e}",
                                "Urdu": f"جواب کی تخلیق ناکام ہو گئی: {e}"
                            }
                            st.error(error_msg.get(detected_lang, f"Response generation failed: {e}"))

if __name__ == "__main__":
    main()
